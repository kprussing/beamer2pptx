"""
beamer2pptx
^^^^^^^^^^^

Utilities for converting a beamer presentation into PowerPoint.
"""

import logging
import os
import pathlib
import re
import subprocess
import tempfile

import pptx

from typing import (
    List,
    Optional,
    Sequence,
    Tuple,
    Union,
)
from os import PathLike

logging.getLogger(__name__).addHandler(logging.NullHandler())


ASPECT_RATIOS = {
    "4:3": (pptx.util.Emu(9144000),
            pptx.util.Emu(6858000)
            ),
    "16:9": (pptx.util.Emu(12192000),
             pptx.util.Emu(6858000)
             )
}
"""The known valid aspect ratios."""


def extract_aspect_ratio(path: Union[PathLike, str],
                         timeout: Optional[float] = None,
                         ) -> str:
    """Extract the aspect ratio from the PDF.

    Parameters
    ----------

    path: path-like
        The path to the slides PDF.
    timeout: float, optional
        The timeout to pass to :func:`subprocess.run`.

    Returns
    -------

    str:
        The aspect ratio of the presentation in the format 'w:h'.

    Raises
    ------

    subprocess.TimeoutError:
        If the call to :manpage:`pdfinfo` times out.
    subprocess.CalledProcessError:
        If the call to :manpage:`pdfinfo` raises an error.
    RuntimeError:
        If the aspect ratio cannot be deduced.

    Notes
    -----

    This calls :manpage:`pdfinfo` to get the page size and compute the
    aspect ratio.  It checks against the known possibilities rounded to
    five significant figures and returns the most appropriate one.

    """
    proc = subprocess.run(["pdfinfo", str(path)],
                          capture_output=True,
                          timeout=timeout,
                          text=True,
                          )

    logger = logging.getLogger(__name__ + ".extract_aspect_ratio")
    try:
        proc.check_returncode()
    except Exception:
        logger.error(f"Error message: '{proc.stderr}'")
        raise

    if proc.stdout == "":
        raise RuntimeError(f"'{logger.name}' no output from pdfinfo")

    match = re.search(r"Page\s*size:\s*([\d.]+)\s*.\s*([\d.]+)",
                      proc.stdout,
                      re.MULTILINE
                      )
    if not match:
        raise RuntimeError(
            f"'{logger.name}' could not locate page size"
        )

    ratio = float(match.group(1)) / float(match.group(2))
    for key, (width, height) in ASPECT_RATIOS.items():
        value = width / height
        if round(value, ndigits=4) == round(ratio, ndigits=4):
            return key

    raise RuntimeError(
        f"'{logger.name}' could not deduce aspect ratio for {ratio}"
    )


def extract_metadata(path: Union[PathLike, str],
                     timeout: Optional[float] = None,
                     ) -> Tuple[str, str, str, str]:
    """Extract the metadata from the PDF.

    Parameters
    ----------

    path: path-like
        The path to the slides PDF.
    timeout: float, optional
        The timeout to pass to :func:`subprocess.run`.

    Returns
    -------

    str:
        The title of the presentation.
    str:
        The subject of the presentation.
    str:
        The keywords in the presentation.
    str:
        The author of the presentation.

    Raises
    ------

    subprocess.TimeoutError:
        If the call to :manpage:`pdfinfo` times out.
    subprocess.CalledProcessError:
        If the call to :manpage:`pdfinfo` raises an error.

    Notes
    -----

    This calls :manpage:`pdfinfo` to extract the metadata.

    """
    proc = subprocess.run(["pdfinfo", str(path)],
                          capture_output=True,
                          timeout=timeout,
                          text=True,
                          )

    logger = logging.getLogger(__name__ + ".extract_metadata")
    try:
        proc.check_returncode()
    except Exception:
        logger.error(f"Error message: '{proc.stderr}'")
        raise

    title = ""
    subject = ""
    keywords = ""
    author = ""
    for line in proc.stdout.splitlines():
        match = re.match(r"Title:\s*(.*)", line, re.IGNORECASE)
        if match:
            title = match.group(1)
            continue

        match = re.match(r"Subject:\s*(.*)", line, re.IGNORECASE)
        if match:
            subject = match.group(1)
            continue

        match = re.match(r"Keywords:\s*(.*)", line, re.IGNORECASE)
        if match:
            keywords = match.group(1)
            continue

        match = re.match(r"Author:\s*(.*)", line, re.IGNORECASE)
        if match:
            author = match.group(1)
            continue

    return title, subject, keywords, author


def extract_notes(path: Union[PathLike, str],
                  timeout: Optional[float] = None,
                  ) -> List[str]:
    """Extract the notes from the PDF.

    Parameters
    ----------

    path: path-like
        The path to the notes PDF.
    timeout: float, optional
        The timeout to pass to :func:`subprocess.run`.

    Returns
    -------

    list of str:
        The text of each note page.

    Raises
    ------

    subprocess.TimeoutError:
        If the call to :manpage:`pdftotext` times out.
    subprocess.CalledProcessError:
        If the call to :manpage:`pdftotext` raises an error.

    Notes
    -----

    This calls :manpage:`pdftotext` to do the work and splits the result
    at the formfeeds.

    """
    proc = subprocess.run(["pdftotext", str(path), "-"],
                          capture_output=True,
                          timeout=timeout,
                          text=True,
                          )
    try:
        proc.check_returncode()
    except Exception:
        logger = logging.getLogger(__name__ + ".extract_notes")
        logger.error(f"Error message: '{proc.stderr}'")
        raise

    # Skip the last one which is just the trailing formfeed.
    return proc.stdout.split("\f")[:-1]


def extract_slides(path: Union[PathLike, str],
                   directory: str = os.curdir,
                   timeout: Optional[float] = None,
                   ) -> List[str]:
    """Extract the slides from the PDF.

    Parameters
    ----------

    path: path-like
        The path to the slides PDF.
    directory: str
        The path to the output directory in which to write the images.
    timeout: float, optional
        The timeout to pass to :func:`subprocess.run`.

    Returns
    -------

    list of str:
        The path to each created slide file.

    Raises
    ------

    subprocess.TimeoutError:
        If the call to :manpage:`pdftocairo` times out.
    subprocess.CalledProcessError:
        If the call to :manpage:`pdftocairo` raises an error.

    Notes
    -----

    This calls :manpage:`pdftocairo` to do the work converting each
    slide page to an image.

    """
    output = pathlib.Path(directory).joinpath(pathlib.Path(path).stem)
    proc = subprocess.run(["pdftocairo", "-png", "-r", "600", "-transp",
                           str(path), str(output)
                           ],
                          capture_output=True,
                          timeout=timeout,
                          text=True,
                          )

    logger = logging.getLogger(__name__ + ".extract_slides")
    try:
        proc.check_returncode()
    except Exception:
        logger.error(f"Error message: '{proc.stderr}'")
        raise

    if proc.stdout != "":
        logger.info(proc.stdout)

    return [str(_) for _ in output.parent.glob(output.stem + "*.png")]


def convert(slides: Union[PathLike, str],
            notes: Optional[Union[PathLike, str]] = None,
            notes_map: Sequence[int] = [],
            timeout: Optional[float] = None,
            ) -> pptx.Presentation:
    """Convert the presentation to PowerPoint.

    Parameters
    ----------

    slides: path-like
        The path to the slides PDF.
    notes: path-like, optional
        The path to the notes PDF.
    notes_map: sequence of integers, optional
        The slides to which to assign the notes.
    timeout: float, optional
        The timeout to pass to subroutines.

    Raises
    ------

    subprocess.TimeoutError:
        If a call to one of the subroutines times out.
    subprocess.CalledProcessError:
        If a call to one of the subroutines errors.
    NotImplementedError:
        If the aspect ratio is unknown.
    ValueError:
        If notes_map is present and not the same length as the notes.

    Notes
    -----

    The ``notes_map`` indicates which slides to which to assign the
    notes in the presentation.  It it is not present, the notes are
    assigned to the slides starting at the beginning and ending when no
    notes are left.  If it is present, it must be the same length as the
    number of notes slides extracted from ``notes``.  Each entry in
    ``notes_map`` is the slide (zero based) to which to add the notes.

    """
    logger = logging.getLogger(f"{__name__}.convert")
    pres = pptx.Presentation()

    title, subject, keywords, author = extract_metadata(slides)
    pres.core_properties.title = title
    pres.core_properties.subject = subject
    pres.core_properties.keywords = keywords
    pres.core_properties.author = author

    logger.info("Adjust the aspect ratio")
    aspect = extract_aspect_ratio(slides)
    if aspect in ASPECT_RATIOS:
        pres.slide_width, pres.slide_height = ASPECT_RATIOS[aspect]
    else:
        raise NotImplementedError(
            f"'{logger.name}' unknown aspect ration {aspect}"
        )

    logger.info("Extract the notes")
    notes_text = [] if notes is None else extract_notes(notes, timeout)
    logger.debug(f"Found {len(notes_text)} notes")
    if len(notes_map) == 0:
        notes_map = range(len(notes_text))

    if len(notes_map) != len(notes_text):
        raise ValueError(
            f"'{logger.name}' incompatible note slides to mapping "
            f"({len(notes_text)} vs. {len(notes_map)})"
        )

    BLANK_SLIDE = pres.slide_layouts[6]
    with tempfile.TemporaryDirectory() as temp:
        logger.info("Generate the slide images")
        images = extract_slides(slides, temp, timeout)

        if len(notes_text) > len(images):
            logger.warn(
                f"More notes found than slides "
                f"({len(notes_text)} vs. {len(images)})"
            )

        for count, image in enumerate(sorted(images)):
            slide = pres.slides.add_slide(BLANK_SLIDE)
            slide.shapes.add_picture(image, 0, 0, width=pres.slide_width)
            if count in notes_map:
                slide.notes_slide.notes_text_frame.text = notes_text[
                    notes_map.index(count)
                ]

    return pres
